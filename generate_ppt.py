from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()

    # Define slide layouts
    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]

    # --- Slide 1: Title Slide ---
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "ScholiFi"
    
    # Customizing subtitle
    tf = subtitle.text_frame
    tf.text = "AI-Driven School Supply Chain\n\n"
    p = tf.add_paragraph()
    p.text = "Team: Code//Blooded\n"
    p.font.bold = True
    
    p2 = tf.add_paragraph()
    p2.text = "Participants: Sohil Chavhan, Shorya Ranjan, Awneesh Kumar, Shivansh Pandey"
    p2.font.size = Pt(16)

    # --- Slide 2: Problem Statement ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "Problem Statement"
    tf = slide.placeholders[1].text_frame
    tf.text = "School procurement is currently plagued by inefficiencies:"
    tf.add_paragraph().text = "Manual, paper-heavy, and opaque processes."
    tf.add_paragraph().text = "Budget leaks and unoptimized spending across departments."
    tf.add_paragraph().text = "Local vendor selection is prone to fraud and lacks competitive bidding."
    tf.add_paragraph().text = "No real-time tracking of department budgets vs. actual spending."

    # --- Slide 3: Proposed Solution ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "Proposed Solution"
    tf = slide.placeholders[1].text_frame
    tf.text = "ScholiFi: An automated ERP system tailored for schools."
    tf.add_paragraph().text = "Centralized dashboard for real-time budget tracking."
    tf.add_paragraph().text = "AI-powered digitization of paper invoices to reduce manual data entry."
    tf.add_paragraph().text = "Transparent vendor bidding portal to democratize local supplier contracts."

    # --- Slide 4: Key Features (Including X-Factor) ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "Key Features"
    tf = slide.placeholders[1].text_frame
    tf.text = "Simulated AI OCR for instant invoice digitization and cross-referencing."
    tf.add_paragraph().text = "Automated Budget-Matching: Blocks purchase orders (POs) if they exceed the department's annual budget."
    tf.add_paragraph().text = "Vendor Bidding Portal: Transparent bidding system for local school contracts."
    
    p = tf.add_paragraph()
    p.text = "The X-Factor: Dynamic, premium \"Light Blue & Gold\" user interface with real-time feedback loops that turns a boring ERP into a consumer-grade experience."
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 102, 204) # Highlight color

    # --- Slide 5: System Architecture ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "System Architecture"
    tf = slide.placeholders[1].text_frame
    tf.text = "Frontend: React (Vite) + Tailwind CSS for a highly responsive, animated UI."
    tf.add_paragraph().text = "State Management: React Hooks (useState, useEffect) for dynamic, localized state simulation."
    tf.add_paragraph().text = "Component Structure: Modular architecture (Dashboard, Scanner, Budgets, Vendor Portal)."
    tf.add_paragraph().text = "Mock Data Layer: Simulated real-time API responses for budgets, POs, and bids."

    # --- Slide 6: Technology Stack ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "Technology Stack"
    tf = slide.placeholders[1].text_frame
    tf.text = "Core Framework: React 19 via Vite for blazing-fast development."
    tf.add_paragraph().text = "Styling: Tailwind CSS v3 for rapid, utility-first premium styling."
    tf.add_paragraph().text = "Iconography: Lucide React for consistent, high-quality vector icons."
    tf.add_paragraph().text = "Build Tooling: Node.js, npm, and PostCSS."

    # --- Slide 7: Future Scope ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "Future Scope"
    tf = slide.placeholders[1].text_frame
    tf.text = "Integration with real OCR models (e.g., Google Cloud Vision or Tesseract) for actual document processing."
    tf.add_paragraph().text = "Full backend architecture (Node.js/Express + PostgreSQL) for real-time, persistent data."
    tf.add_paragraph().text = "Machine Learning models for predictive budget forecasting and anomaly detection."
    tf.add_paragraph().text = "Mobile application for quick PO approvals by principals and administrators."

    # --- Slide 8: Working Prototype ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "Working Prototype"
    tf = slide.placeholders[1].text_frame
    tf.text = "The prototype is fully functional and live locally."
    tf.add_paragraph().text = "URL: http://localhost:5173"
    tf.add_paragraph().text = "Demonstrates the Dashboard, simulated Invoice Scanning workflow, Department Budget alerts, and Vendor Bidding Portal."

    # Save presentation
    output_path = "C:\\Users\\Sohil\\OneDrive\\Desktop\\ScholiFi_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully to {output_path}")

if __name__ == "__main__":
    create_presentation()
